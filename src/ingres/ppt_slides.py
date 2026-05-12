import os
import logging
from pathlib import Path
from typing import Any, Dict, List, Optional

from fsspec import AbstractFileSystem
from llama_index.core.readers.base import BaseReader
from llama_index.core.schema import Document
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from typing import IO, cast

log = logging.getLogger(__name__)


class PptxSlideReader(BaseReader):
    """PowerPoint reader. One Document per slide."""

    def load_data(
        self,
        file: str | Path,
        extra_info: Optional[Dict] = None,
        fs: Optional[AbstractFileSystem] = None,
    ) -> List[Document]:
        file = Path(file)

        if fs:
            with fs.open(str(file), "rb") as f:
                presentation = Presentation(cast(IO[bytes], f))
        else:
            presentation = Presentation(str(file))

        documents: List[Document] = []

        for i, slide in enumerate(presentation.slides):
            slide_lines: List[str] = [f"Slide {i + 1}"]

            for shape in slide.shapes:
                # cast to Any — Pylance sees BaseShape but runtime
                # subclasses have has_text_frame / has_table / shape_type
                s: Any = shape

                if s.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    slide_lines.append("[Image]")

                elif s.has_table:
                    for row in s.table.rows:
                        row_text = " | ".join(
                            cell.text.strip()
                            for cell in row.cells
                            if cell.text.strip()
                        )
                        if row_text:
                            slide_lines.append(row_text)

                elif s.has_text_frame:
                    for paragraph in s.text_frame.paragraphs:
                        text = "".join(
                            run.text for run in paragraph.runs
                        ).strip()
                        if text:
                            slide_lines.append(text)

            slide_text = "\n".join(slide_lines)
            if not slide_text.strip():
                continue

            metadata: Dict = dict(extra_info or {})
            metadata.update({
                "level":    1,
                "h1":       f"Slide {i + 1}",
                "h2":       "",
                "h3":       "",
                "category": "PPT",
            })

            documents.append(
                Document(text=slide_text, metadata=metadata)
            )

        return documents